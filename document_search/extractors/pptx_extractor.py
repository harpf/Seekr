from __future__ import annotations

import os
from pathlib import Path

from pptx import Presentation

from document_search.extractors.base import TextExtractor
from document_search.models import ContentBlock, ExtractionResult
from document_search.services.ocr_service import ocr_office_embedded_images


class PptxTextExtractor(TextExtractor):
    def extract(self, file_path: Path) -> ExtractionResult:
        try:
            prs = Presentation(str(file_path))
            blocks: list[ContentBlock] = []
            for i, slide in enumerate(prs.slides, start=1):
                texts: list[str] = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        texts.append(shape.text.strip())
                    if hasattr(shape, "has_table") and shape.has_table:
                        table = shape.table
                        rows = ["\t".join(cell.text.strip() for cell in row.cells) for row in table.rows]
                        texts.append("\n".join(r for r in rows if r.strip()))
                    if hasattr(shape, "alternative_text") and shape.alternative_text:
                        texts.append(f"AltText: {shape.alternative_text.strip()}")
                notes = slide.notes_slide.notes_text_frame.text.strip() if slide.has_notes_slide else ""
                if notes:
                    texts.append(f"Speaker notes: {notes}")
                full = "\n".join([t for t in texts if t])
                if full:
                    blocks.append(ContentBlock("slide", i, full, self.__class__.__name__, {"slide": i}))
            if os.getenv("DOCUMENT_SEARCH_OCR_ENABLED", "false").lower() == "true":
                for i, txt in enumerate(ocr_office_embedded_images(file_path), start=1):
                    blocks.append(ContentBlock("ocr_image", i, txt, self.__class__.__name__, {"source": "embedded_image"}))
            cp = prs.core_properties
            meta = {"title": cp.title, "author": cp.author, "created": str(cp.created) if cp.created else None}
            return ExtractionResult(file_path=file_path, status="ok", document_metadata={"slide_count": len(prs.slides), "pptx_properties": meta}, blocks=blocks)
        except Exception as ex:
            return ExtractionResult(file_path=file_path, status="error", error_message=str(ex))
